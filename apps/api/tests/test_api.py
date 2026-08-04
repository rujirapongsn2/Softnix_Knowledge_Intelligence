import os
import tempfile
from io import BytesIO
from pathlib import Path
from types import SimpleNamespace

_TEST_ROOT = tempfile.mkdtemp()
os.environ["DATABASE_URL"] = f"sqlite:///{_TEST_ROOT}/skip.db"
os.environ["FILE_STORAGE_PATH"] = f"{_TEST_ROOT}/files"
os.environ["INITIAL_ADMIN_PASSWORD"] = "correct-horse-battery-staple"
os.environ["LIGHTRAG_BASE_URL"] = ""
os.environ["REDIS_URL"] = ""
os.environ["OPENROUTER_API_KEY"] = ""
os.environ["EXT_OCR_KEY"] = ""

from fastapi.testclient import TestClient
import httpx
from openpyxl import Workbook
from pptx import Presentation
from app import services
from app.config import Settings
from app.db import SessionLocal
from app.external_ocr import ExternalOcrClient
from app.main import app
from app.models import Document, DocumentChunk, KnowledgeBase, LegalInstrument, LegalInstrumentRelation


def client():
    with TestClient(app) as test_client:
        assert test_client.post("/api/v1/auth/login", json={"username": "admin", "password": "correct-horse-battery-staple"}).status_code == 200
        yield test_client


def test_vertical_slice():
    test_client = next(client())
    kb = test_client.post("/api/v1/knowledge-bases", json={"name": "IT", "code": "enterprise-it"}).json()
    test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/activate")
    uploaded = test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/documents", files={"file": ("architecture.txt", b"Customer Portal runs on APP-01.", "text/plain")}).json()
    assert uploaded["status"] == "queued"
    assert test_client.post("/api/v1/internal/process-next").json()["processed"] is True
    assert test_client.get(f"/api/v1/documents/{uploaded['document_id']}/text").json()["status"] == "completed"
    token = test_client.post("/api/v1/tokens", json={"name": "agent", "allowed_knowledge_base_ids": [kb["id"]], "allowed_tools": ["search_knowledge"]}).json()
    reply = test_client.post("/mcp", headers={"Authorization": f"Bearer {token['token']}"}, json={"jsonrpc": "2.0", "id": 1, "method": "tools/call", "params": {"name": "search_knowledge", "arguments": {"query": "What runs on APP-01?"}}})
    assert reply.status_code == 200
    assert reply.json()["result"]["structuredContent"]["sources"]
    assert reply.json()["result"]["structuredContent"]["request_id"] == 1
    assert reply.json()["result"]["structuredContent"]["metadata"]["retrieval_plan"]["intent"] == "entity_lookup"
    activity = test_client.get("/api/v1/mcp/activity").json()
    call = next(row for row in activity if row["metadata"].get("tool") == "search_knowledge")
    assert call["metadata"]["query"] == "What runs on APP-01?"
    assert any(step["channel"] == "full_text" and step["status"] == "used" for step in call["metadata"]["route"])
    transactions = test_client.get("/api/v1/logs/transactions?limit=100").json()
    mcp_transaction = next(row for row in transactions if row["path"] == "/mcp" and row["retrieval"])
    assert mcp_transaction["retrieval"]["transport"] == "mcp"
    assert mcp_transaction["retrieval"]["retrieval_plan"]["intent"] == "entity_lookup"
    assert any(step["channel"] == "full_text" for step in mcp_transaction["retrieval"]["retrieval_trace"])


def test_document_inventory_summary_is_deterministic_and_search_fallback_preserves_original_query():
    test_client = next(client())
    kb = test_client.post("/api/v1/knowledge-bases", json={"name": "Inventory", "code": "inventory-kb"}).json()
    assert test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/activate").status_code == 200
    first = test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/documents", files={"file": ("one.txt", b"one", "text/plain")}).json()
    second = test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/documents", files={"file": ("two.txt", b"two", "text/plain")}, data={"document_type": "legal"}).json()
    assert first["status"] == second["status"] == "queued"
    token = test_client.post("/api/v1/tokens", json={
        "name": "inventory-agent", "allowed_knowledge_base_ids": [kb["id"]],
        "allowed_tools": ["search_knowledge", "document_inventory_summary"],
    }).json()
    headers = {"Authorization": f"Bearer {token['token']}"}
    listed = test_client.post("/mcp", headers=headers, json={"jsonrpc": "2.0", "id": 1, "method": "tools/list"}).json()
    assert "document_inventory_summary" in {tool["name"] for tool in listed["result"]["tools"]}

    original = "ชุดเอกสารนี้มีกฎหมายทั้งหมดกี่ฉบับ และแบ่งเป็นประเภทใดบ้าง?"
    inventory = test_client.post("/mcp", headers=headers, json={"jsonrpc": "2.0", "id": 2, "method": "tools/call", "params": {
        "name": "document_inventory_summary", "arguments": {"query": original, "scope": "all", "include_documents": True},
    }}).json()["result"]["structuredContent"]
    assert inventory["total_documents"] == 2
    assert {group["key"] for group in inventory["groups"] if group["dimension"] == "document_type"} == {"general", "legal"}
    assert inventory["metadata"]["retrieval_plan"]["intent"] == "document_inventory"
    assert inventory["metadata"]["retrieval_trace"][0]["channel"] == "document_registry"

    fallback = test_client.post("/mcp", headers=headers, json={"jsonrpc": "2.0", "id": 3, "method": "tools/call", "params": {
        "name": "search_knowledge", "arguments": {"query": original},
    }}).json()["result"]["structuredContent"]
    assert fallback["total_documents"] == 2
    assert fallback["metadata"]["retrieval_plan"]["intent"] == "document_inventory"
    activity = test_client.get("/api/v1/mcp/activity").json()
    call = next(row for row in activity if row["metadata"].get("request_id") == "3")
    assert call["metadata"]["query"] == original
    inventory_call = next(row for row in activity if row["metadata"].get("request_id") == "2")
    assert inventory_call["metadata"]["query"] == original
    assert services.is_document_inventory_query("แสดงรายการเอกสารทั้งหมด") is True
    assert services.is_document_inventory_query("ลบเอกสารทั้งหมดอย่างไร") is False
    assert services.is_document_inventory_query("How do I archive all documents?") is False
    # Keep the shared test worker queue isolated for subsequent API tests.
    for _ in range(6):
        if not test_client.post("/api/v1/internal/process-next").json().get("processed"):
            break


def test_documents_page_is_bounded_and_reports_global_processing_state():
    test_client = next(client())
    kb = test_client.post("/api/v1/knowledge-bases", json={"name": "Paged documents", "code": "paged-documents"}).json()
    assert test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/activate").status_code == 200
    first = test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/documents", files={"file": ("land-law.txt", "มาตรา 4".encode("utf-8"), "text/plain")}, data={"document_type": "legal"}).json()
    second = test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/documents", files={"file": ("notes.txt", b"notes", "text/plain")}).json()
    assert first["status"] == second["status"] == "queued"

    page = test_client.get(f"/api/v1/knowledge-bases/{kb['id']}/documents/page?limit=1&search=land-law").json()
    assert page["total"] == 1
    assert len(page["items"]) == 1
    assert page["has_legal_documents"] is True
    assert page["has_completed_documents"] is False
    assert page["processing_count"] == 2
    assert page["items"][0]["processing_job_status"] == "queued"
    assert page["items"][0]["processing_job_progress_percent"] == 0
    # Do not leave initial or legal follow-up jobs in the shared fixture queue.
    processed = []
    for _ in range(5):
        processed.append(test_client.post("/api/v1/internal/process-next").json()["processed"])
    assert any(processed)


def test_mcp_legal_tools_are_agent_safe_and_scope_bound():
    test_client = next(client())
    kb = test_client.post("/api/v1/knowledge-bases", json={"name": "Legal MCP", "code": "legal-mcp-agent"}).json()
    assert test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/activate").status_code == 200
    worklist = test_client.get(f"/api/v1/knowledge-bases/{kb['id']}/legal-registry/worklist")
    assert worklist.status_code == 200 and worklist.json()["ready"] is True
    token = test_client.post("/api/v1/tokens", json={
        "name": "legal-agent", "allowed_knowledge_base_ids": [kb["id"]],
        "allowed_tools": ["resolve_legal_context", "get_legal_instrument", "get_provision_history"],
    }).json()
    headers = {"Authorization": f"Bearer {token['token']}"}
    listed = test_client.post("/mcp", headers=headers, json={"jsonrpc": "2.0", "id": 1, "method": "tools/list"}).json()
    names = {tool["name"] for tool in listed["result"]["tools"]}
    assert {"resolve_legal_context", "get_legal_instrument", "get_provision_history"}.issubset(names)
    resolved = test_client.post("/mcp", headers=headers, json={"jsonrpc": "2.0", "id": 2, "method": "tools/call", "params": {
        "name": "resolve_legal_context", "arguments": {"query": "มาตรา 39", "knowledge_base_ids": ["client-cannot-expand-scope"]}
    }}).json()
    structured = resolved["result"]["structuredContent"]
    assert structured["knowledge_base_ids"] == [kb["id"]]
    assert structured["metadata"]["source_of_truth"] == "PostgreSQL legal registry"
    denied = test_client.post("/mcp", headers=headers, json={"jsonrpc": "2.0", "id": 3, "method": "tools/call", "params": {
        "name": "get_legal_instrument", "arguments": {"instrument_id": "outside-scope"}
    }}).json()
    assert denied["error"]["code"] == "LEGAL_INSTRUMENT_NOT_FOUND"
    # A rejected MCP call must still be visible in the Trace Explorer -- not
    # just the raw audit log -- since record_retrieval_execution only ever
    # fires on the success path.
    error_traces = test_client.get("/api/v1/traces?tool=get_legal_instrument&status=error").json()
    assert any(item["status"] == "error" for item in error_traces)


def test_knowledge_base_can_be_disabled_activated_and_safely_deleted():
    test_client = next(client())
    empty = test_client.post("/api/v1/knowledge-bases", json={"name": "Disposable", "code": "disposable-kb"}).json()
    assert test_client.post(f"/api/v1/knowledge-bases/{empty['id']}/disable").status_code == 200
    assert test_client.delete(f"/api/v1/knowledge-bases/{empty['id']}").json()["status"] == "deleted"
    assert empty["id"] not in {kb["id"] for kb in test_client.get("/api/v1/knowledge-bases").json()}

    kb = test_client.post("/api/v1/knowledge-bases", json={"name": "Managed", "code": "managed-kb"}).json()
    assert test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/disable").status_code == 200
    rejected = test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/documents", files={"file": ("notes.txt", b"blocked", "text/plain")})
    assert rejected.status_code == 409 and rejected.json()["error"]["code"] == "KNOWLEDGE_BASE_DISABLED"
    assert test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/activate").status_code == 200
    uploaded = test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/documents", files={"file": ("notes.txt", b"allowed", "text/plain")}).json()
    nonempty = test_client.delete(f"/api/v1/knowledge-bases/{kb['id']}")
    assert nonempty.status_code == 409 and nonempty.json()["error"]["code"] == "KNOWLEDGE_BASE_NOT_EMPTY"
    assert uploaded["status"] == "queued"
    assert test_client.post("/api/v1/internal/process-next").json()["processed"] is True


def test_knowledge_base_code_is_derived_safely_for_non_latin_and_repeated_names():
    test_client = next(client())
    first = test_client.post("/api/v1/knowledge-bases", json={"name": "กฎหมายคุ้มครองข้อมูลส่วนบุคคล"})
    second = test_client.post("/api/v1/knowledge-bases", json={"name": "กฎหมายคุ้มครองข้อมูลส่วนบุคคล"})
    assert first.status_code == 200 and second.status_code == 200
    assert first.json()["code"].startswith("kb-")
    assert second.json()["code"] != first.json()["code"]

    duplicate = test_client.post("/api/v1/knowledge-bases", json={"name": "Another label", "code": first.json()["code"]})
    assert duplicate.status_code == 409
    assert duplicate.json()["error"]["code"] == "KNOWLEDGE_BASE_CODE_EXISTS"


def test_retrieval_policy_can_be_updated_and_is_returned_in_kb_contract():
    test_client = next(client())
    kb = test_client.post("/api/v1/knowledge-bases", json={"name": "Planner policy", "code": "planner-policy"}).json()
    updated = test_client.patch(f"/api/v1/knowledge-bases/{kb['id']}/retrieval-config", json={
        "retrieval_mode": "precision", "enable_lightrag": False, "maximum_graph_depth": 2,
    })
    assert updated.status_code == 200
    assert updated.json()["retrieval_config"]["retrieval_mode"] == "precision"
    assert updated.json()["retrieval_config"]["enable_lightrag"] is False
    assert updated.json()["retrieval_config"]["maximum_graph_depth"] == 2


def test_legal_document_type_queues_automatic_metadata_extraction():
    test_client = next(client())
    kb = test_client.post("/api/v1/knowledge-bases", json={"name": "Automatic legal", "code": "automatic-legal-kb"}).json()
    uploaded = test_client.post(
        f"/api/v1/knowledge-bases/{kb['id']}/documents",
        files={"file": ("contract.txt", b"Party A shall provide support.", "text/plain")},
        data={"document_type": "contract"},
    )
    assert uploaded.status_code == 200
    assert uploaded.json()["document_type"] == "contract"
    assert uploaded.json()["legal_extraction_automatic"] is True
    assert test_client.post("/api/v1/internal/process-next").json()["processed"] is True
    document = test_client.get(f"/api/v1/documents/{uploaded.json()['document_id']}/text").json()
    assert document["status"] == "completed"
    assert document["document_type"] == "contract"
    jobs = test_client.get(f"/api/v1/documents/{uploaded.json()['document_id']}/jobs").json()
    assert any(job["type"] == "EXTRACT_LEGAL_METADATA" and job["status"] == "queued" for job in jobs)
    # Drain the follow-up job so independent tests do not consume it.
    assert test_client.post("/api/v1/internal/process-next").json()["processed"] is True


def test_upload_rejects_unknown_document_type():
    test_client = next(client())
    kb = test_client.post("/api/v1/knowledge-bases", json={"name": "Typed", "code": "invalid-type-kb"}).json()
    response = test_client.post(
        f"/api/v1/knowledge-bases/{kb['id']}/documents",
        files={"file": ("notes.txt", b"notes", "text/plain")},
        data={"document_type": "unknown"},
    )
    assert response.status_code == 400
    assert response.json()["error"]["code"] == "DOCUMENT_TYPE_INVALID"


def test_token_is_not_returned_after_creation():
    test_client = next(client())
    token = test_client.post("/api/v1/tokens", json={"name": "agent", "allowed_tools": ["search_knowledge"]}).json()
    assert token["token"].startswith("skik_live_")
    listed = test_client.post("/mcp", headers={"Authorization": f"Bearer {token['token']}"}, json={"jsonrpc": "2.0", "id": 1, "method": "tools/list"}).json()
    assert "tools" in listed["result"]


def test_refresh_cookie_restores_access_session():
    test_client = next(client())
    assert test_client.cookies.get("skip_refresh")
    test_client.cookies.delete("skip_access")
    refreshed = test_client.post("/api/v1/auth/refresh")
    assert refreshed.status_code == 200
    assert test_client.cookies.get("skip_access")
    assert test_client.get("/api/v1/auth/me").status_code == 200


def test_auto_retrieval_fixture_exercises_scopes_exact_dates_and_rerank_policy(monkeypatch):
    test_client = next(client())
    kb = test_client.post("/api/v1/knowledge-bases", json={"name": "Auto retrieval fixture", "code": "auto-retrieval-fixture"}).json()
    test_client.patch(f"/api/v1/knowledge-bases/{kb['id']}/retrieval-config", json={
        "enable_lightrag": False, "planner_llm_fallback": False, "enable_reranker": False,
    })
    from datetime import date, datetime
    import hashlib
    with SessionLocal() as db:
        documents = []
        for index, (name, text, published_at) in enumerate([
            ("vpn.txt", "ขั้นตอนขอสิทธิ์ VPN คือ submit request and manager approval.", None),
            ("delay.txt", "ปัจจัยหลักที่ทำให้โครงการล่าช้า คือ vendor dependency.", None),
            ("abc-june.txt", "ข่าวบริษัท ABC เดือนมิถุนายน 2026 เปิดศูนย์บริการ.", date(2026, 6, 15)),
            ("abc-may.txt", "ข่าวบริษัท ABC เดือนพฤษภาคม 2026 ไม่ควรถูกคืน.", date(2026, 5, 15)),
            ("SNX-2026-001.txt", "เอกสารเลขที่ SNX-2026-001: Information Security Standard.", None),
            ("land-transfer.txt", "ขั้นตอนการโอนกรรมสิทธิ์หรือสิทธิครอบครองในที่ดินที่มีโฉนดที่ดินหรือหนังสือรับรองการทำประโยชน์ ต้องยื่นคำขอต่อสำนักงานที่ดินและแสดงเอกสารสิทธิ์.", None),
        ]):
            digest = hashlib.sha256(f"{name}:{text}".encode()).hexdigest()
            doc = Document(knowledge_base_id=kb["id"], original_filename=name, stored_filename=name,
                           storage_path=f"/tmp/{name}", mime_type="text/plain", file_size=len(text), checksum_sha256=digest,
                           title=name.removesuffix(".txt"), document_type="general", published_at=published_at,
                           tags=[], status="completed", extracted_text=text, indexed_at=datetime.utcnow())
            db.add(doc); db.flush()
            db.add(DocumentChunk(document_id=doc.id, knowledge_base_id=kb["id"], chunk_index=0, content=text,
                                 content_sha256=hashlib.sha256(text.encode()).hexdigest(), char_start=0, char_end=len(text), token_count=len(text.split())))
            documents.append(doc)
        db.commit()
    transfer_document_id = next(document.id for document in documents if document.original_filename == "land-transfer.txt")
    from app.retrieval import RetrievalEvidence
    import time
    original_vector = services.query_database_vectors

    def vector_probe(db, query, kb_ids, limit, trace=None, plan=None):
        if query != "การโอนกรรมสิทธิ์หรือสิทธิครอบครองในที่ดินที่มีโฉนดที่ดินหรือหนังสือรับรองการทำประโยชน์ ต้องทำอย่างไร?":
            return original_vector(db, query, kb_ids, limit, trace, plan)
        services._append_retrieval_trace(trace, channel="semantic_vector", system="test vector probe", status="used",
                                         started_at=time.monotonic(), result_count=1, detail="fixture embedding")
        return RetrievalEvidence([{"citation_id": "V1", "document_id": transfer_document_id,
                                   "chunk_id": "fixture-vector-chunk", "chunk_index": 0,
                                   "title": "land-transfer", "excerpt": "fixture vector evidence", "relevance": 1.0}], [], [], [])

    monkeypatch.setattr(services, "query_database_vectors", vector_probe)
    architecture = documents[0]
    app_entity = test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/entities", json={"name": "APP-01", "entity_type": "Application"}).json()
    portal = test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/entities", json={"name": "Customer Portal", "entity_type": "Application"}).json()
    relationship = test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/relationships", json={
        "source_entity_id": app_entity["id"], "target_entity_id": portal["id"], "relationship_type": "CONNECTS_TO",
        "document_id": architecture.id, "excerpt": "APP-01 connects to Customer Portal.",
    })
    assert relationship.status_code == 200
    cases = {
        "ขั้นตอนขอสิทธิ์ VPN คืออะไร": ["vector", "full_text"],
        "APP-01 เชื่อมต่อกับระบบใด": ["graph"],
        "ระบบใดได้รับผลกระทบหาก APP-01 ล่ม": ["graph", "vector"],
        "ปัจจัยหลักที่ทำให้โครงการล่าช้า": ["graph", "vector"],
        "ข่าวบริษัท ABC เดือนมิถุนายน 2026": ["full_text", "vector"],
        "เอกสารเลขที่ SNX-2026-001": ["exact_document", "full_text"],
        "ภาพรวมความสัมพันธ์ระหว่างหน่วยงาน": ["graph", "vector"],
        "การโอนกรรมสิทธิ์หรือสิทธิครอบครองในที่ดินที่มีโฉนดที่ดินหรือหนังสือรับรองการทำประโยชน์ ต้องทำอย่างไร?": ["vector", "full_text"],
    }
    results = {}
    for query, channels in cases.items():
        result = test_client.post("/api/v1/query", json={"query": query, "knowledge_base_ids": [kb["id"]]}).json()
        assert [item for item in result["metadata"]["retrieval_plan"]["channels"]] == channels
        results[query] = result
    local_trace = results["APP-01 เชื่อมต่อกับระบบใด"]["metadata"]["retrieval_trace"]
    assert next(item for item in local_trace if item["channel"] == "graph")["result_count"] == 1
    impact_trace = results["ระบบใดได้รับผลกระทบหาก APP-01 ล่ม"]["metadata"]["retrieval_trace"]
    assert next(item for item in impact_trace if item["channel"] == "graph")["result_count"] == 1
    global_trace = results["ภาพรวมความสัมพันธ์ระหว่างหน่วยงาน"]["metadata"]["retrieval_trace"]
    assert next(item for item in global_trace if item["channel"] == "graph")["result_count"] == 1
    news_sources = results["ข่าวบริษัท ABC เดือนมิถุนายน 2026"]["sources"]
    assert news_sources and {source["title"] for source in news_sources} == {"abc-june"}
    exact_trace = results["เอกสารเลขที่ SNX-2026-001"]["metadata"]["retrieval_trace"]
    assert next(item for item in exact_trace if item["channel"] == "exact_document")["result_count"] == 1
    assert all(item["channel"] != "graph" or item["status"] == "skipped" for item in exact_trace)
    transfer_trace = results["การโอนกรรมสิทธิ์หรือสิทธิครอบครองในที่ดินที่มีโฉนดที่ดินหรือหนังสือรับรองการทำประโยชน์ ต้องทำอย่างไร?"]["metadata"]["retrieval_trace"]
    assert next(item for item in transfer_trace if item["channel"] == "semantic_vector")["status"] == "used"
    assert next(item for item in transfer_trace if item["channel"] == "full_text")["status"] == "used"
    assert results["การโอนกรรมสิทธิ์หรือสิทธิครอบครองในที่ดินที่มีโฉนดที่ดินหรือหนังสือรับรองการทำประโยชน์ ต้องทำอย่างไร?"]["sources"]
    rerank_trace = results["เอกสารเลขที่ SNX-2026-001"]["metadata"]["retrieval_trace"]
    assert next(item for item in rerank_trace if item["channel"] == "rerank")["detail"] == "disabled by retrieval policy"
    authority = test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/entities", json={"name": "กรมทดสอบ", "entity_type": "Organization"}).json()
    department = test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/entities", json={"name": "ฝ่ายทดสอบ", "entity_type": "Organization"}).json()
    assert test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/relationships", json={
        "source_entity_id": authority["id"], "target_entity_id": department["id"], "relationship_type": "OVERSEES",
        "document_id": documents[2].id, "excerpt": "กรมทดสอบดูแลฝ่ายทดสอบ",
    }).status_code == 200
    filtered_global = test_client.post("/api/v1/query", json={
        "knowledge_base_ids": [kb["id"]], "query": "ภาพรวมความสัมพันธ์ระหว่างหน่วยงาน",
        "filters": {"published_from": "2026-06-01", "published_to": "2026-07-01"},
    }).json()
    assert any(source["document_id"] == documents[2].id for source in filtered_global["sources"])
    assert all(source["document_id"] != documents[3].id for source in filtered_global["sources"])


def test_revoked_token_cannot_call_mcp():
    test_client = next(client())
    token = test_client.post("/api/v1/tokens", json={"name": "agent", "allowed_tools": ["search_knowledge"]}).json()
    assert test_client.post(f"/api/v1/tokens/{token['id']}/revoke").status_code == 200
    response = test_client.post("/mcp", headers={"Authorization": f"Bearer {token['token']}"}, json={"jsonrpc": "2.0", "id": 1, "method": "tools/list"})
    assert response.status_code == 200
    assert response.json()["error"]["code"] == "AUTH_TOKEN_REVOKED"


def test_lightrag_honors_publication_filter_before_fusion():
    test_client = next(client())
    kb = test_client.post("/api/v1/knowledge-bases", json={"name": "LightRAG filter", "code": "lightrag-filter"}).json()
    from datetime import date, datetime
    import hashlib
    with SessionLocal() as db:
        documents = []
        for name, published_at in (("in-range.txt", date(2026, 6, 10)), ("out-of-range.txt", date(2026, 5, 10))):
            text = f"LightRAG source {name}"
            document = Document(knowledge_base_id=kb["id"], original_filename=name, stored_filename=name,
                                storage_path=f"/tmp/{name}", mime_type="text/plain", file_size=len(text),
                                checksum_sha256=hashlib.sha256(f"{kb['id']}:{name}".encode()).hexdigest(),
                                title=name.removesuffix(".txt"), document_type="general", published_at=published_at,
                                tags=[], status="completed", extracted_text=text, indexed_at=datetime.utcnow())
            db.add(document); db.flush(); documents.append(document)
        db.commit()
        from app.planner import RetrievalPlan
        fake_engine = SimpleNamespace(query=lambda *_: services.RetrievalEvidence([
            {"citation_id": "S1", "document_id": documents[0].id, "title": "in-range", "chunk_id": "1", "excerpt": "in", "relevance": 1.0},
            {"citation_id": "S2", "document_id": documents[1].id, "title": "out-of-range", "chunk_id": "2", "excerpt": "out", "relevance": 0.9},
        ], [], [], []))
        plan = RetrievalPlan(intent="news_by_date", channels=["lightrag"], published_from=date(2026, 6, 1), published_to=date(2026, 7, 1))
        evidence = services._query_lightrag(db, fake_engine, "news", [kb["id"]], 10, [], plan)
    assert [source["document_id"] for source in evidence.sources] == [documents[0].id]


def test_token_rotation_replaces_secret_without_exposing_old_key():
    test_client = next(client())
    original = test_client.post("/api/v1/tokens", json={"name": "rotate-me", "allowed_tools": ["search_knowledge"]}).json()
    rotated = test_client.post(f"/api/v1/tokens/{original['id']}/rotate")
    assert rotated.status_code == 200
    replacement = rotated.json()
    assert replacement["token"] != original["token"]
    assert replacement["name"] == original["name"]
    assert replacement["allowed_tools"] == original["allowed_tools"]
    old_response = test_client.post("/mcp", headers={"Authorization": f"Bearer {original['token']}"}, json={"jsonrpc": "2.0", "id": 1, "method": "tools/list"})
    assert old_response.json()["error"]["code"] == "AUTH_TOKEN_REVOKED"
    new_response = test_client.post("/mcp", headers={"Authorization": f"Bearer {replacement['token']}"}, json={"jsonrpc": "2.0", "id": 2, "method": "tools/list"})
    assert new_response.json()["result"]["tools"]


def test_token_rotation_preserves_disabled_state():
    test_client = next(client())
    original = test_client.post("/api/v1/tokens", json={"name": "disabled-rotation", "allowed_tools": ["search_knowledge"]}).json()
    assert test_client.post(f"/api/v1/tokens/{original['id']}/disable").status_code == 200
    rotated = test_client.post(f"/api/v1/tokens/{original['id']}/rotate")
    assert rotated.status_code == 200
    replacement = rotated.json()
    assert replacement["status"] == "inactive"
    blocked = test_client.post("/mcp", headers={"Authorization": f"Bearer {replacement['token']}"}, json={"jsonrpc": "2.0", "id": 3, "method": "tools/list"})
    assert blocked.json()["error"]["code"] == "AUTH_TOKEN_INVALID"


def test_batch_upload_queues_each_file_with_shared_document_type():
    test_client = next(client())
    kb = test_client.post("/api/v1/knowledge-bases", json={"name": "Batch upload"}).json()
    response = test_client.post(
        f"/api/v1/knowledge-bases/{kb['id']}/documents/batch",
        files=[
            ("files", ("act.txt", "มาตรา 1 applies.".encode(), "text/plain")),
            ("files", ("notice.txt", "ประกาศแก้ไข มาตรา 2".encode(), "text/plain")),
        ],
        data={"document_type": "legal"},
    )
    assert response.status_code == 200
    payload = response.json()
    assert payload["status"] == "queued"
    assert payload["queued_count"] == 2
    assert [item["document_type"] for item in payload["results"]] == ["legal", "legal"]
    assert all(item["status"] == "queued" and item["job_id"] for item in payload["results"])
    for _ in range(3):
        test_client.post("/api/v1/internal/process-next")


def test_batch_upload_isolates_duplicate_failure_from_other_files():
    test_client = next(client())
    kb = test_client.post("/api/v1/knowledge-bases", json={"name": "Batch partial"}).json()
    assert test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/documents", files={"file": ("existing.txt", b"existing", "text/plain")}).status_code == 200
    response = test_client.post(
        f"/api/v1/knowledge-bases/{kb['id']}/documents/batch",
        files=[
            ("files", ("existing.txt", b"existing", "text/plain")),
            ("files", ("new.txt", b"new independent content", "text/plain")),
        ],
        data={"document_type": "general"},
    )
    assert response.status_code == 200
    payload = response.json()
    assert payload["status"] == "partial"
    assert payload["queued_count"] == 1 and payload["failed_count"] == 1
    assert next(item for item in payload["results"] if item["filename"] == "existing.txt")["error_code"] == "FILE_DUPLICATE"
    assert next(item for item in payload["results"] if item["filename"] == "new.txt")["status"] == "queued"
    for _ in range(3):
        test_client.post("/api/v1/internal/process-next")


def test_document_list_exposes_latest_follow_up_job_status():
    test_client = next(client())
    kb = test_client.post("/api/v1/knowledge-bases", json={"name": "Legal job status"}).json()
    uploaded = test_client.post(
        f"/api/v1/knowledge-bases/{kb['id']}/documents",
        files={"file": ("law.txt", "มาตรา 1 applies.".encode(), "text/plain")},
        data={"document_type": "legal"},
    ).json()
    rows = test_client.get(f"/api/v1/knowledge-bases/{kb['id']}/documents").json()
    row = next(item for item in rows if item["id"] == uploaded["document_id"])
    assert row["processing_job_status"] == "queued"
    assert row["processing_job_type"] == "PROCESS_DOCUMENT"
    for _ in range(3):
        test_client.post("/api/v1/internal/process-next")


def test_mcp_only_reads_active_knowledge_bases():
    test_client = next(client())
    active = test_client.post("/api/v1/knowledge-bases", json={"name": "MCP active", "code": "mcp-active-kb"}).json()
    disabled = test_client.post("/api/v1/knowledge-bases", json={"name": "MCP disabled", "code": "mcp-disabled-kb"}).json()
    test_client.post(f"/api/v1/knowledge-bases/{active['id']}/activate")
    test_client.post(f"/api/v1/knowledge-bases/{disabled['id']}/activate")
    active_doc = test_client.post(f"/api/v1/knowledge-bases/{active['id']}/documents", files={"file": ("active.txt", b"Active MCP evidence.", "text/plain")}).json()
    assert test_client.post("/api/v1/internal/process-next").json()["processed"] is True
    token = test_client.post("/api/v1/tokens", json={"name": "active-only-agent", "allowed_knowledge_base_ids": [active["id"], disabled["id"]], "allowed_tools": ["search_knowledge"]}).json()
    assert test_client.post(f"/api/v1/knowledge-bases/{disabled['id']}/disable").status_code == 200

    # The client-supplied KB value is ignored; the token scope is authoritative.
    active_response = test_client.post("/mcp", headers={"Authorization": f"Bearer {token['token']}"}, json={"jsonrpc": "2.0", "id": 1, "method": "tools/call", "params": {"name": "search_knowledge", "arguments": {"query": "Active MCP evidence", "knowledge_base_ids": ["client-controlled-value"]}}})
    assert active_response.json()["result"]["structuredContent"]["sources"][0]["document_id"] == active_doc["document_id"]
    call = next(row for row in test_client.get("/api/v1/mcp/activity").json() if row["metadata"].get("request_id") == "1")
    assert call["metadata"]["knowledge_base_ids"] == [active["id"]]
    assert call["metadata"]["client_knowledge_base_ids_ignored"] is True
    disabled_response = test_client.post("/mcp", headers={"Authorization": f"Bearer {token['token']}"}, json={"jsonrpc": "2.0", "id": 2, "method": "tools/call", "params": {"name": "search_knowledge", "arguments": {"query": "anything", "knowledge_base_ids": [disabled["id"]]}}})
    assert disabled_response.json()["result"]["structuredContent"]["metadata"]["knowledge_base_ids"] == [active["id"]]
    assert test_client.post(f"/api/v1/knowledge-bases/{active['id']}/disable").status_code == 200
    no_active_scope = test_client.post("/mcp", headers={"Authorization": f"Bearer {token['token']}"}, json={"jsonrpc": "2.0", "id": 3, "method": "tools/call", "params": {"name": "search_knowledge", "arguments": {"query": "anything"}}})
    assert no_active_scope.json()["error"]["code"] == "KNOWLEDGE_BASE_INACTIVE"


def test_mcp_token_scope_must_reference_active_knowledge_bases():
    test_client = next(client())
    draft = test_client.post("/api/v1/knowledge-bases", json={"name": "Draft MCP scope", "code": "draft-mcp-scope"}).json()
    response = test_client.post("/api/v1/tokens", json={
        "name": "draft-scoped-agent",
        "allowed_knowledge_base_ids": [draft["id"]],
        "allowed_tools": ["search_knowledge"],
    })
    assert response.status_code == 400
    assert response.json()["error"]["code"] == "KNOWLEDGE_BASE_INACTIVE"


def test_graph_impact_returns_direct_indirect_impacts_and_sources():
    test_client = next(client())
    kb = test_client.post("/api/v1/knowledge-bases", json={"name": "Graph KB", "code": "graph-kb"}).json()
    uploaded = test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/documents", files={"file": ("architecture.txt", b"Architecture evidence", "text/plain")}).json()
    test_client.post("/api/v1/internal/process-next")
    app = test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/entities", json={"name": "APP-01", "entity_type": "Application", "document_id": uploaded["document_id"], "excerpt": "APP-01 hosts the portal."}).json()
    portal = test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/entities", json={"name": "Customer Portal", "entity_type": "Application"}).json()
    process = test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/entities", json={"name": "Customer Service", "entity_type": "BusinessProcess"}).json()
    test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/relationships", json={"source_entity_id": portal["id"], "target_entity_id": app["id"], "relationship_type": "RUNS_ON", "document_id": uploaded["document_id"], "excerpt": "Customer Portal runs on APP-01."})
    test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/relationships", json={"source_entity_id": process["id"], "target_entity_id": portal["id"], "relationship_type": "USES", "document_id": uploaded["document_id"], "excerpt": "Customer Service uses Customer Portal."})
    result = test_client.post("/api/v1/query/impact", json={"subject": "APP-01", "scenario": "stops", "knowledge_base_ids": [kb["id"]], "max_depth": 2}).json()
    assert result["subject"]["name"] == "APP-01"
    assert [item["name"] for item in result["direct_impacts"]] == ["Customer Portal"]
    assert [item["name"] for item in result["indirect_impacts"]] == ["Customer Service"]
    assert result["sources"][0]["citation_id"] == "S1"
    assert test_client.get("/api/v1/system/graph-projection").json()["events"]["queued"] >= 3


def test_document_list_preview_and_reprocess():
    test_client = next(client())
    kb = test_client.post("/api/v1/knowledge-bases", json={"name": "Docs", "code": "docs-kb"}).json()
    uploaded = test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/documents", files={"file": ("notes.txt", b"The previewable document text.", "text/plain")}).json()
    assert len(test_client.get(f"/api/v1/knowledge-bases/{kb['id']}/documents").json()) == 1
    test_client.post("/api/v1/internal/process-next")
    assert "previewable" in test_client.get(f"/api/v1/documents/{uploaded['document_id']}/text").json()["text"]
    reprocess = test_client.post(f"/api/v1/documents/{uploaded['document_id']}/reprocess").json()
    assert reprocess["status"] == "queued"
    duplicate = test_client.post(f"/api/v1/documents/{uploaded['document_id']}/reprocess")
    assert duplicate.status_code == 409
    assert duplicate.json()["error"]["code"] == "DOCUMENT_PROCESSING_IN_PROGRESS"
    assert len(test_client.get(f"/api/v1/documents/{uploaded['document_id']}/jobs").json()) == 2


def test_markitdown_extracts_structured_office_markdown_and_legacy_fallback(tmp_path, monkeypatch):
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Platform Architecture"
    text_box = slide.shapes.add_textbox(100, 100, 1000, 500)
    text_box.text_frame.text = "APP-01 runs the customer portal."
    pptx_path = tmp_path / "architecture.pptx"
    presentation.save(pptx_path)
    pptx_markdown = services.extract_text(SimpleNamespace(storage_path=str(pptx_path)))
    assert "# Platform Architecture" in pptx_markdown
    assert "APP-01 runs the customer portal." in pptx_markdown

    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Systems"
    worksheet.append(["System", "Owner"])
    worksheet.append(["APP-01", "Platform Team"])
    xlsx_path = tmp_path / "systems.xlsx"
    workbook.save(xlsx_path)
    xlsx_markdown = services.extract_text(SimpleNamespace(storage_path=str(xlsx_path)))
    assert "## Systems" in xlsx_markdown and "| APP-01 | Platform Team |" in xlsx_markdown

    text_path = tmp_path / "fallback.txt"
    text_path.write_text("Legacy parser keeps this extraction available.", encoding="utf-8")
    monkeypatch.setattr(services, "_markitdown_extract", lambda _: (_ for _ in ()).throw(RuntimeError("converter failed")))
    assert "Legacy parser" in services.extract_text(SimpleNamespace(storage_path=str(text_path)))

    scanned_pdf = tmp_path / "scanned.pdf"
    scanned_pdf.write_bytes(b"not a text-layer PDF")
    monkeypatch.setattr(services, "_markitdown_extract", lambda _: "")
    try:
        services.extract_text(SimpleNamespace(storage_path=str(scanned_pdf)))
    except RuntimeError as error:
        assert str(error) == "OCR_REQUIRED"
    else:
        raise AssertionError("scanned PDFs must require OCR")


def test_external_ocr_v3_returns_markdown_and_reports_progress():
    calls, progress = [], []

    def handler(request: httpx.Request):
        calls.append((request.method, request.url.path))
        if request.method == "POST":
            assert b'disable_structure' in request.content
            return httpx.Response(200, json={"job_id": "ocr-job-1"})
        if request.url.path.endswith("/status") and len([call for call in calls if call[1].endswith("/status")]) == 1:
            return httpx.Response(200, json={"status": "processing", "progress": {"percent": 50, "stage": "ocr"}})
        if request.url.path.endswith("/status"):
            return httpx.Response(200, json={"status": "completed"})
        if request.url.path.endswith("/result"):
            return httpx.Response(200, json={"results": {"combined_markdown": "# Scanned policy\n\nArticle 1"}})
        raise AssertionError(request.url.path)

    settings = Settings(ext_ocr_key="test-key", ext_ocr_base_url="https://ocr.test", ext_ocr_poll_interval_seconds=0)
    client = httpx.Client(transport=httpx.MockTransport(handler), base_url="https://ocr.test")
    result = ExternalOcrClient(settings, client).extract_markdown(Path(__file__), lambda stage, percent: progress.append((stage, percent)))
    assert result == "# Scanned policy\n\nArticle 1"
    assert progress[0] == ("external_ocr_queued", 15)
    assert any(stage == "external_ocr_ocr" for stage, _ in progress)


def test_pptx_and_xlsx_uploads_are_processed_as_markdown():
    test_client = next(client())
    kb = test_client.post("/api/v1/knowledge-bases", json={"name": "Office", "code": "office-kb"}).json()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Service Catalog"
    pptx_stream = BytesIO(); presentation.save(pptx_stream)
    pptx = test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/documents", files={"file": ("catalog.pptx", pptx_stream.getvalue(), "application/vnd.openxmlformats-officedocument.presentationml.presentation")}).json()
    for _ in range(10):
        test_client.post("/api/v1/internal/process-next")
        pptx_preview = test_client.get(f"/api/v1/documents/{pptx['document_id']}/text").json()
        if pptx_preview["status"] != "queued":
            break
    assert pptx_preview["status"] == "completed", pptx_preview
    assert "# Service Catalog" in pptx_preview["text"]

    workbook = Workbook(); worksheet = workbook.active; worksheet.title = "Inventory"; worksheet.append(["System", "Status"]); worksheet.append(["APP-01", "Active"])
    xlsx_stream = BytesIO(); workbook.save(xlsx_stream)
    xlsx = test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/documents", files={"file": ("inventory.xlsx", xlsx_stream.getvalue(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")}).json()
    for _ in range(10):
        test_client.post("/api/v1/internal/process-next")
        xlsx_preview = test_client.get(f"/api/v1/documents/{xlsx['document_id']}/text").json()
        if xlsx_preview["status"] != "queued":
            break
    assert xlsx_preview["status"] == "completed", xlsx_preview
    assert "| APP-01 | Active |" in xlsx_preview["text"]


def test_upload_rejects_mime_type_that_does_not_match_extension():
    test_client = next(client())
    kb = test_client.post("/api/v1/knowledge-bases", json={"name": "Mime", "code": "mime-kb"}).json()
    response = test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/documents", files={"file": ("unsafe.xlsx", b"not-an-excel-file", "text/plain")})
    assert response.status_code == 400
    assert response.json()["error"]["code"] == "FILE_MIME_TYPE_NOT_SUPPORTED"


def test_legal_metadata_articles_amendments_crud():
    test_client = next(client())
    kb = test_client.post("/api/v1/knowledge-bases", json={"name": "Legal", "code": "legal-kb"}).json()
    uploaded = test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/documents", files={"file": ("law.txt", "มาตรา 1 applies.".encode(), "text/plain")}, data={"document_type": "legal"}).json()
    test_client.post("/api/v1/internal/process-next")
    metadata = {"document_type": "ประกาศ", "articles": [{"article_number": "1", "text": "applies", "evidence_quote": "มาตรา 1 applies"}], "amendments": [{"title": "ประกาศแก้ไข", "announcement_number": "ฉบับที่ 2", "changes": "แก้ไขมาตรา 1", "evidence_quote": "ฉบับที่ 2"}]}
    saved = test_client.put(f"/api/v1/documents/{uploaded['document_id']}/legal-metadata", json={"metadata": metadata})
    assert saved.status_code == 200 and saved.json()["legal_metadata"]["articles"][0]["article_number"] == "1"
    patched = test_client.patch(f"/api/v1/documents/{uploaded['document_id']}/legal-metadata", json={"metadata": {"confidence": 0.9}})
    assert patched.json()["legal_metadata"]["confidence"] == 0.9
    graph_sync = test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/graph/sync").json()
    assert graph_sync["entities"] >= 3 and graph_sync["relationships"] >= 2
    assert test_client.get(f"/api/v1/documents/{uploaded['document_id']}/text").json()["legal_metadata"]["amendments"][0]["title"] == "ประกาศแก้ไข"
    assert test_client.delete(f"/api/v1/documents/{uploaded['document_id']}/legal-metadata").json()["status"] == "deleted"
    assert test_client.get(f"/api/v1/documents/{uploaded['document_id']}/text").json()["legal_metadata"] is None


def test_legal_graph_v2_keeps_provisions_document_scoped_and_reviews_suggestions():
    test_client = next(client())
    kb = test_client.post("/api/v1/knowledge-bases", json={"name": "Legal graph v2", "code": "legal-graph-v2"}).json()
    first_metadata = {
        "schema_version": 2,
        "instrument": {"kind": "Act", "official_title": "PDPA Act", "official_number": "2562"},
        "provisions": [{"kind": "article", "number": "1", "evidence_quote": "มาตรา 1 แห่งพระราชบัญญัตินี้"}],
        "parties": [], "obligations": [], "rights": [], "prohibitions": [], "penalties": [], "definitions": [], "amendments": [], "references": [],
    }
    second_metadata = {
        "schema_version": 2,
        "instrument": {"kind": "Notification", "official_title": "PDPC Security Notification", "official_number": "1/2565"},
        "provisions": [{"kind": "article", "number": "1", "evidence_quote": "ข้อ 1 ของประกาศนี้"}],
        "parties": [], "obligations": [], "rights": [], "prohibitions": [], "penalties": [], "definitions": [], "amendments": [],
        "references": [{"relationship": "ISSUED_UNDER", "target_title": "PDPA Act", "target_number": "2562", "evidence_quote": "อาศัยอำนาจตาม PDPA Act", "confidence": 0.9}],
    }
    with SessionLocal() as db:
        db.add_all([
            Document(knowledge_base_id=kb["id"], original_filename="act.txt", stored_filename="act.txt", storage_path="/tmp/act.txt", mime_type="text/plain", file_size=1,
                     checksum_sha256="a" * 64, title="PDPA Act", document_type="legal", status="completed", extracted_text="มาตรา 1 แห่งพระราชบัญญัตินี้", legal_metadata=first_metadata),
            Document(knowledge_base_id=kb["id"], original_filename="notice.txt", stored_filename="notice.txt", storage_path="/tmp/notice.txt", mime_type="text/plain", file_size=1,
                     checksum_sha256="b" * 64, title="PDPC Security Notification", document_type="regulation", status="completed", extracted_text="อาศัยอำนาจตาม PDPA Act", legal_metadata=second_metadata),
        ])
        db.commit()
        first_build = services.rebuild_legal_graph(db, kb["id"])
        assert first_build["documents"] == 2

    verified = test_client.get(f"/api/v1/knowledge-bases/{kb['id']}/legal-graph?view=verified").json()
    provisions = [node for node in verified["nodes"] if node["entity_type"] == "Provision"]
    assert len(provisions) == 2
    assert len({node["identity_key"] for node in provisions}) == 2
    suggested = test_client.get(f"/api/v1/knowledge-bases/{kb['id']}/legal-graph?view=suggested").json()
    edge = next(edge for edge in suggested["edges"] if edge["relationship_type"] == "ISSUED_UNDER")
    assert edge["review_status"] == "suggested" and edge["sources"][0]["excerpt"] == "อาศัยอำนาจตาม PDPA Act"
    approved = test_client.patch(f"/api/v1/relationships/{edge['id']}/legal-review", json={"status": "verified", "note": "validated"})
    assert approved.status_code == 200 and approved.json()["review_status"] == "verified"
    with SessionLocal() as db:
        services.rebuild_legal_graph(db, kb["id"])
    verified_after = test_client.get(f"/api/v1/knowledge-bases/{kb['id']}/legal-graph?view=verified").json()
    assert any(item["id"] == edge["id"] and item["review_status"] == "verified" for item in verified_after["edges"])
    definition = next(node for node in verified["nodes"] if node["entity_type"] == "Provision")
    inspector = test_client.get(f"/api/v1/entities/{definition['id']}/inspector?depth=1")
    assert inspector.status_code == 200
    payload = inspector.json()
    assert payload["entity"]["entity_type"] == "Provision"
    assert payload["evidence"] and payload["context"]["documents"]
    assert payload["relationships"]["outgoing"] or payload["relationships"]["incoming"]
    assert "versions" in payload and "warnings" in payload["analysis"]
    legal_map = test_client.get(f"/api/v1/knowledge-bases/{kb['id']}/legal-map?view=verified").json()
    assert legal_map["mode"] == "map"
    assert len(legal_map["instruments"]) == 2
    assert all("entity_count" in item and "relationship_count" in item for item in legal_map["instruments"])
    summary = legal_map["relationship_summary"]
    assert {"verified", "suggested", "rejected", "manual", "internal", "cross_document"} <= summary.keys()
    assert summary["internal"] >= 2 and summary["cross_document"] >= 0
    with SessionLocal() as db:
        source_instrument = db.query(LegalInstrument).filter_by(knowledge_base_id=kb["id"]).first()
        db.add(LegalInstrumentRelation(
            knowledge_base_id=kb["id"], source_instrument_id=source_instrument.id,
            relation="REFERS_TO", target_text="Unresolved instrument", origin="ai_suggestion", review_status="suggested",
        ))
        db.commit()
    summary_with_unresolved = test_client.get(f"/api/v1/knowledge-bases/{kb['id']}/legal-map?view=verified").json()["relationship_summary"]
    assert summary_with_unresolved["cross_document"] == summary["cross_document"] + 1
    assert summary_with_unresolved["suggested"] == summary["suggested"] + 1
    instrument_view = test_client.get(f"/api/v1/knowledge-bases/{kb['id']}/legal-map?view=verified&instrument_id={legal_map['instruments'][0]['id']}").json()
    assert instrument_view["mode"] == "instrument"
    assert instrument_view["nodes"]
    assert all(node["entity_type"] in {"LegalInstrument", "Provision"} for node in instrument_view["nodes"])
    assert instrument_view["edges"]
    assert all("sources" in edge for edge in instrument_view["edges"])


def test_database_chunk_retrieval_is_scoped_to_knowledge_base():
    test_client = next(client())
    first = test_client.post("/api/v1/knowledge-bases", json={"name": "First", "code": "first-kb"}).json()
    second = test_client.post("/api/v1/knowledge-bases", json={"name": "Second", "code": "second-kb"}).json()
    test_client.post(f"/api/v1/knowledge-bases/{first['id']}/documents", files={"file": ("first.txt", b"APP-01 owns a unique primary platform.", "text/plain")})
    test_client.post(f"/api/v1/knowledge-bases/{second['id']}/documents", files={"file": ("second.txt", b"APP-01 owns a unique secondary platform.", "text/plain")})
    assert test_client.post("/api/v1/internal/process-next").json()["processed"] is True
    assert test_client.post("/api/v1/internal/process-next").json()["processed"] is True
    result = test_client.post("/api/v1/query", json={"knowledge_base_ids": [first["id"]], "query": "unique platform"}).json()
    assert result["sources"]
    assert all(source["document_id"] != "" for source in result["sources"])
    assert all("primary" in source["excerpt"] for source in result["sources"])


def test_embedding_reindex_queues_completed_documents_without_full_reprocessing():
    test_client = next(client())
    kb = test_client.post("/api/v1/knowledge-bases", json={"name": "Reindex", "code": "reindex-kb"}).json()
    uploaded = test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/documents", files={"file": ("notes.txt", b"Reindex this document.", "text/plain")}).json()
    for _ in range(10):
        test_client.post("/api/v1/internal/process-next")
        if test_client.get(f"/api/v1/documents/{uploaded['document_id']}/text").json()["status"] == "completed":
            break
    queued = test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/documents/reindex?force=true").json()
    assert queued == {"status": "queued", "count": 1}
    for _ in range(10):
        test_client.post("/api/v1/internal/process-next")
        jobs = test_client.get(f"/api/v1/documents/{uploaded['document_id']}/jobs").json()
        if jobs[0]["status"] == "completed":
            break
    jobs = test_client.get(f"/api/v1/documents/{uploaded['document_id']}/jobs").json()
    assert jobs[0]["status"] == "completed"
    assert jobs[0]["stage"] == "completed"


def test_metrics_and_audit_log_are_exposed():
    test_client = next(client())
    kb = test_client.post("/api/v1/knowledge-bases", json={"name": "Audit", "code": "audit-kb"}).json()
    audit = test_client.get("/api/v1/audit-logs").json()
    assert any(row["action"] == "knowledge_base.create" and row["target_id"] == kb["id"] for row in audit)
    metrics = test_client.get("/metrics")
    assert metrics.status_code == 200
    assert "softnix_http_requests_total" in metrics.text


def test_request_transactions_are_operator_safe_and_include_request_metadata():
    test_client = next(client())
    created = test_client.post("/api/v1/knowledge-bases", json={"name": "Transaction", "code": "transaction-kb"})
    request_id = created.headers["X-Request-ID"]
    transactions = test_client.get("/api/v1/logs/transactions?limit=100")
    assert transactions.status_code == 200
    row = next(item for item in transactions.json() if item["request_id"] == request_id)
    assert row["method"] == "POST"
    assert row["path"] == "/api/v1/knowledge-bases"
    assert row["status_code"] == 200
    assert isinstance(row["duration_ms"], (int, float))
    assert row["authentication"] == "admin_session"


def test_admin_query_transaction_includes_retrieval_executor_trace():
    test_client = next(client())
    kb = test_client.post("/api/v1/knowledge-bases", json={"name": "Trace", "code": "trace-kb"}).json()
    test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/documents", files={"file": ("trace.txt", b"The trace document describes the platform.", "text/plain")})
    assert test_client.post("/api/v1/internal/process-next").json()["processed"] is True
    response = test_client.post("/api/v1/query", json={"query": "platform", "knowledge_base_ids": [kb["id"]]})
    assert response.status_code == 200
    request_id = response.headers["X-Request-ID"]
    transactions = test_client.get("/api/v1/logs/transactions?limit=100").json()
    row = next(item for item in transactions if item["request_id"] == request_id)
    assert row["retrieval"]["transport"] == "api"
    assert row["retrieval"]["retrieval_plan"]["intent"]
    assert row["retrieval"]["retrieval_trace"]
    traces = test_client.get("/api/v1/traces?transport=api").json()
    trace = next(item for item in traces if item["request_id"] == request_id)
    detail = test_client.get(f"/api/v1/traces/{trace['trace_id']}").json()
    assert detail["root_span"]["span_id"] == "root"
    assert detail["spans"]
    assert all("offset_ms" in span and "duration_ms" in span for span in detail["spans"])
    assert detail["request_summary"]["query_preview"] == "platform"
    assert detail["request_summary"]["query_sha256"]
    assert "Bearer " not in str(detail)
    paged = test_client.get("/api/v1/traces?transport=api&limit=1&paginate=true").json()
    assert "items" in paged and isinstance(paged["items"], list)
    assert paged["items"][0]["trace_id"] == trace["trace_id"]
    transaction_page = test_client.get("/api/v1/logs/transactions?limit=1&paginate=true").json()
    assert "items" in transaction_page and transaction_page["items"]


def test_document_restore_graph_layout_and_feedback_lifecycle():
    test_client = next(client())
    kb = test_client.post("/api/v1/knowledge-bases", json={"name": "Lifecycle", "code": "lifecycle-kb"}).json()
    uploaded = test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/documents", files={"file": ("notes.txt", b"Lifecycle evidence for APP-02.", "text/plain")}).json()
    document_id = uploaded["document_id"]
    assert test_client.delete(f"/api/v1/documents/{document_id}").json()["status"] == "deleted"
    assert test_client.get(f"/api/v1/knowledge-bases/{kb['id']}/documents").json() == []
    assert test_client.post(f"/api/v1/documents/{document_id}/restore").json()["status"] == "queued"
    assert test_client.get(f"/api/v1/knowledge-bases/{kb['id']}/documents?include_deleted=true").json()[0]["deleted_at"] is None
    entity = test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/entities", json={"name": "APP-02", "entity_type": "Application"}).json()
    assert test_client.put(f"/api/v1/knowledge-bases/{kb['id']}/graph-layout", json={"items": [{"entity_id": entity["id"], "x": 120, "y": 80}]}).json()["status"] == "success"
    assert test_client.get(f"/api/v1/knowledge-bases/{kb['id']}/graph-layout").json()["items"][0]["x"] == 120
    assert test_client.patch(f"/api/v1/entities/{entity['id']}", json={"name": "APP-02 API"}).json()["name"] == "APP-02 API"
    test_client.post("/api/v1/internal/process-next")
    result = test_client.post("/api/v1/query", json={"knowledge_base_ids": [kb["id"]], "query": "APP-02"}).json()
    feedback = test_client.post(f"/api/v1/query/results/{result['result_id']}/feedback", json={"rating": 1, "comment": "useful"})
    assert feedback.status_code == 200
    assert feedback.json()["status"] == "success"


def test_mcp_rate_limit_returns_jsonrpc_error():
    test_client = next(client())
    token = test_client.post("/api/v1/tokens", json={
        "name": "limited", "requests_per_minute": 1, "allowed_tools": ["search_knowledge"],
    }).json()
    request = {"jsonrpc": "2.0", "id": 7, "method": "tools/list"}
    assert test_client.post("/mcp", headers={"Authorization": f"Bearer {token['token']}"}, json=request).json()["result"]["tools"]
    response = test_client.post("/mcp", headers={"Authorization": f"Bearer {token['token']}"}, json=request)
    assert response.status_code == 200
    assert response.json()["error"]["code"] == "MCP_RATE_LIMITED"


def test_custom_document_template_keeps_a_safe_processing_profile_and_user_metadata():
    test_client = next(client())
    kb = test_client.post("/api/v1/knowledge-bases", json={"name": "Template KB", "code": "template-kb"}).json()
    created = test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/document-templates", json={
        "name": "Official notification", "base_document_type": "regulation",
        "fields": [
            {"key": "issuer", "label": "Issuing organisation", "field_type": "text", "required": True,
             "filterable": True, "graph_entity_type": "Organization", "graph_relationship": "ISSUED_BY"},
            {"key": "effective_date", "label": "Effective date", "field_type": "date", "filterable": True},
            {"key": "internal_note", "label": "Internal note", "field_type": "text", "searchable": False},
        ],
    })
    assert created.status_code == 200
    template = created.json()
    assert template["base_document_type"] == "regulation"
    assert template["is_system"] is False
    template_fields = {field["key"]: field for field in template["fields"]}
    assert "reference_number" in template_fields  # profile baseline is inherited
    assert template_fields["issuer"]["filterable"] is True
    assert template_fields["issuer"]["graph_relationship"] == "ISSUED_BY"
    listed = test_client.get(f"/api/v1/knowledge-bases/{kb['id']}/document-templates").json()
    assert {"General document", "Official notification"}.issubset({row["name"] for row in listed})
    uploaded = test_client.post(
        f"/api/v1/knowledge-bases/{kb['id']}/documents",
        files={"file": ("notification.txt", "ประกาศ เรื่อง ทดสอบ".encode("utf-8"), "text/plain")},
        data={"template_id": template["id"], "metadata_json": '{"issuer":"สำนักงานทดสอบ","effective_date":"2026-01-01","internal_note":"ไม่ควรค้นพบ"}'},
    )
    assert uploaded.status_code == 200
    assert uploaded.json()["document_type"] == "regulation"
    preview = test_client.get(f"/api/v1/documents/{uploaded.json()['document_id']}/text").json()
    assert preview["metadata_template_id"] == template["id"]
    assert preview["document_metadata"]["issuer"] == "สำนักงานทดสอบ"
    assert preview["document_metadata"]["internal_note"] == "ไม่ควรค้นพบ"
    filtered_page = test_client.get(f"/api/v1/knowledge-bases/{kb['id']}/documents/page", params={"template_id": template["id"]})
    assert filtered_page.status_code == 200
    assert filtered_page.json()["total"] == 1
    invalid = test_client.post(
        f"/api/v1/knowledge-bases/{kb['id']}/documents",
        files={"file": ("invalid.txt", b"test", "text/plain")},
        data={"template_id": template["id"], "metadata_json": '{"unknown":"value"}'},
    )
    assert invalid.status_code == 400
    updated = test_client.patch(f"/api/v1/documents/{uploaded.json()['document_id']}/metadata", json={"values": {"issuer": "สำนักงานใหม่", "effective_date": "2026-01-01", "internal_note": "หมายเหตุลับ"}})
    assert updated.status_code == 200
    assert updated.json()["document_metadata"] == {"issuer": "สำนักงานใหม่", "effective_date": "2026-01-01", "internal_note": "หมายเหตุลับ"}
    assert test_client.post("/api/v1/internal/process-next").json()["processed"] is True
    filtered_query = test_client.post("/api/v1/query", json={
        "knowledge_base_ids": [kb["id"]], "query": "ประกาศ", "filters": {"metadata": {"issuer": "สำนักงานใหม่", "effective_date": "2026-01-01"}},
    })
    assert filtered_query.status_code == 200
    assert filtered_query.json()["sources"]
    assert {source["document_id"] for source in filtered_query.json()["sources"]} == {uploaded.json()["document_id"]}
    assert filtered_query.json()["metadata"]["retrieval_plan"]["metadata_filters"] == {"issuer": "สำนักงานใหม่", "effective_date": "2026-01-01"}
    unsearchable_query = test_client.post("/api/v1/query", json={
        "knowledge_base_ids": [kb["id"]], "query": "หมายเหตุลับ", "filters": {"metadata": {"issuer": "สำนักงานใหม่", "effective_date": "2026-01-01"}},
    })
    assert unsearchable_query.status_code == 200
    assert not unsearchable_query.json()["sources"]
    entities = test_client.get(f"/api/v1/knowledge-bases/{kb['id']}/entities?search=สำนักงานใหม่").json()
    assert any(entity["name"] == "สำนักงานใหม่" for entity in entities)
    graph_entity = next(entity for entity in entities if entity["name"] == "สำนักงานใหม่")
    graph = test_client.get(f"/api/v1/entities/{graph_entity['id']}/graph").json()
    assert any(edge["type"] == "ISSUED_BY" for edge in graph["edges"])
    retired = test_client.delete(f"/api/v1/document-templates/{template['id']}")
    assert retired.status_code == 200
    listed_all = test_client.get(f"/api/v1/knowledge-bases/{kb['id']}/document-templates?include_inactive=true")
    assert listed_all.status_code == 200
    listed_by_name = {row["name"]: row for row in listed_all.json()}
    assert {"General document", "Legal document", "Official notification"}.issubset(listed_by_name)
    assert listed_by_name["Official notification"]["is_active"] is False
    assert listed_by_name["Official notification"]["usage_count"] == 1
    assert listed_by_name["General document"]["is_system"] is True
    # Retiring a template blocks future uploads but does not erase its
    # document-scoped metadata schema or prevent correcting historic data.
    after_retirement = test_client.patch(
        f"/api/v1/documents/{uploaded.json()['document_id']}/metadata",
        json={"values": {"issuer": "สำนักงานหลังปิดประเภท"}},
    )
    assert after_retirement.status_code == 200
    assert after_retirement.json()["document_metadata"] == {"issuer": "สำนักงานหลังปิดประเภท"}
    blocked = test_client.post(
        f"/api/v1/knowledge-bases/{kb['id']}/documents",
        files={"file": ("blocked.txt", b"test", "text/plain")},
        data={"template_id": template["id"], "metadata_json": '{"issuer":"x"}'},
    )
    assert blocked.status_code == 400
    activated = test_client.post(f"/api/v1/document-templates/{template['id']}/activate")
    assert activated.status_code == 200
    select_invalid = test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/document-templates", json={
        "name": "Invalid select", "fields": [{"key": "kind", "label": "Kind", "field_type": "select"}],
    })
    assert select_invalid.status_code == 422


def test_document_template_api_rejects_duplicate_fields_and_does_not_keep_old_profile_defaults():
    test_client = next(client())
    kb = test_client.post("/api/v1/knowledge-bases", json={"name": "Template validation", "code": "template-validation"}).json()
    duplicate = test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/document-templates", json={
        "name": "Duplicate fields", "base_document_type": "regulation",
        "fields": [
            {"key": "owner", "label": "Owner"},
            {"key": "owner", "label": "Owner again"},
        ],
    })
    assert duplicate.status_code == 422

    created = test_client.post(f"/api/v1/knowledge-bases/{kb['id']}/document-templates", json={
        "name": "Switchable type", "base_document_type": "regulation",
        "fields": [{"key": "owner", "label": "Owner"}],
    }).json()
    switched = test_client.patch(f"/api/v1/document-templates/{created['id']}", json={"base_document_type": "general"})
    assert switched.status_code == 200
    assert {field["key"] for field in switched.json()["fields"]} == {"owner"}
def test_document_template_guards_system_names_and_archived_empty_schema_snapshots():
    test_client = next(client())
    kb = test_client.post("/api/v1/knowledge-bases", json={"name": "Template Guard KB", "code": "template-guard-kb"}).json()

    duplicate_system_name = test_client.post(
        f"/api/v1/knowledge-bases/{kb['id']}/document-templates",
        json={"name": "Legal document"},
    )
    assert duplicate_system_name.status_code == 409

    created = test_client.post(
        f"/api/v1/knowledge-bases/{kb['id']}/document-templates",
        json={"name": "Empty metadata type"},
    )
    assert created.status_code == 200
    template = created.json()

    uploaded = test_client.post(
        f"/api/v1/knowledge-bases/{kb['id']}/documents",
        files={"file": ("empty-schema.txt", b"test", "text/plain")},
        data={"template_id": template["id"], "metadata_json": "{}"},
    )
    assert uploaded.status_code == 200
    document_id = uploaded.json()["document_id"]

    retired = test_client.delete(f"/api/v1/document-templates/{template['id']}")
    assert retired.status_code == 200
    after_retirement = test_client.patch(
        f"/api/v1/documents/{document_id}/metadata",
        json={"values": {}},
    )
    assert after_retirement.status_code == 200

    active_override = test_client.patch(
        f"/api/v1/document-templates/{template['id']}",
        json={"is_active": True},
    )
    assert active_override.status_code == 422


def test_sync_document_metadata_graph_counts_every_new_entity_it_creates():
    next(client())  # ensure the schema exists (created on app startup)
    with SessionLocal() as db:
        kb = KnowledgeBase(code="metadata-graph-count", name="Metadata graph count")
        db.add(kb); db.flush()
        doc = Document(knowledge_base_id=kb.id, original_filename="doc.txt", stored_filename="doc.txt",
                       storage_path="/tmp/doc.txt", mime_type="text/plain", file_size=1,
                       checksum_sha256="ha" * 32, title="Metadata doc", document_type="general", status="completed",
                       metadata_template_fields=[
                           {"key": "issuer", "label": "Issuer", "graph_entity_type": "Organization", "graph_relationship": "ISSUED_BY"},
                           {"key": "recipient", "label": "Recipient", "graph_entity_type": "Organization", "graph_relationship": "SENT_TO"},
                       ],
                       document_metadata={"issuer": "Ministry A", "recipient": "Ministry B"})
        db.add(doc); db.commit()
        # 1 anchor "Document" entity + 2 mapped-field target entities (issuer,
        # recipient) must all be counted -- previously only the anchor was.
        result = services.sync_document_metadata_graph(db, doc)
        assert result["entities"] == 3
        assert result["relationships"] == 2
